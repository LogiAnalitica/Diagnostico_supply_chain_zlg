"""
generar_datos.py  —  v2.0  (arquitectura multi-módulo)
=======================================================
Genera las tres salidas del diagnóstico a partir de los scores calculados
por scoring.py:

    1. dashboard_output.html          — dashboard interactivo para el cliente
    2. presentacion_diagnostico.pptx  — presentación con los gráficos
    3. Matriz Priorización.html       — matriz de priorización de estrategias

Cambios respecto a v1.0
-----------------------
- La sección "DATOS DE ENTRADA" manual fue eliminada.
  Los datos ahora llegan desde scoring.py (calculados a partir de las
  respuestas del cliente en Streamlit).
- construir_payload() ahora recibe los scores como parámetro en lugar
  de leer variables globales.
- Los radares por dimensión se construyen dinámicamente desde CONFIG_MODULOS,
  eliminando las claves hardcodeadas (radar_tecnica, radar_kpi, etc.).
  Cada módulo declara sus propias dimensiones y configuración de radar.
- generar_pptx() itera sobre radares_dimensiones en lugar de buscar
  claves fijas, haciéndolo compatible con cualquier módulo futuro.
- Se mantiene compatibilidad hacia atrás: las claves originales
  (radar_tecnica, radar_kpi, radar_estructura) se conservan como alias
  en el payload para que el dashboard_template.html actual no requiera
  ningún cambio.

Para agregar un módulo nuevo (ej. Transporte)
---------------------------------------------
1. Agrega su entrada en CONFIG_MODULOS con sus dimensiones y benchmarks.
2. Agrega sus benchmarks en BENCHMARKS.
3. Nada más. El resto del código es genérico.
"""

import json
import base64
import os


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 1 — CONFIGURACIÓN POR MÓDULO
# Aquí se define la identidad visual y estructural de cada módulo.
# Agregar un nuevo módulo = agregar una nueva entrada en este diccionario.
# ══════════════════════════════════════════════════════════════════════════════

CONFIG_MODULOS = {

    "MOD-01": {
        "nombre"          : "Almacenamiento",
        "centro_evaluado" : "Almacenamiento",

        # Cada dimensión del módulo define:
        #   clave        → nombre de la variable en el payload (mantiene compatibilidad HTML)
        #   id_dim       → ID de la dimensión en el Excel de parametrización
        #   titulo       → título que aparece en el gráfico radar y en el PPTX
        #   gridshape    → "circular" (radares con muchos ejes) o "linear" (pocos ejes)
        #   es_estructura→ True solo cuando el radar usa configuración especial de eje radial
        "radares_dimensiones": [
            {
                "clave"        : "radar_tecnica",
                "id_dim"       : "DIM-01",
                "titulo"       : "Dimensión Técnica y Operacional",
                "gridshape"    : "circular",
                "es_estructura": False,
            },
            {
                "clave"        : "radar_estructura",
                "id_dim"       : "DIM-02",
                "titulo"       : "Dimensión Estructura",
                "gridshape"    : "linear",
                "es_estructura": True,
            },
            {
                "clave"        : "radar_kpi",
                "id_dim"       : "DIM-03",
                "titulo"       : "Dimensión KPI",
                "gridshape"    : "circular",
                "es_estructura": False,
            },
        ],

        # Layout de la diapositiva 3 del PPTX:
        #   radar_principal   → ocupa la franja superior derecha (más grande)
        #   radares_inferiores→ se reparten en la franja inferior derecha
        "layout_pptx": {
            "radar_principal"   : "radar_tecnica",
            "radares_inferiores": ["radar_kpi", "radar_estructura"],
        },
    },

    # ── Módulos futuros ───────────────────────────────────────────────────────
    # Ejemplo de cómo se verá MOD-04 (Planeación) cuando esté listo:
    #
    # "MOD-04": {
    #     "nombre"         : "Planeación",
    #     "centro_evaluado": "Planeación de la demanda",
    #     "radares_dimensiones": [
    #         {"clave": "radar_demanda",    "id_dim": "DIM-01",
    #          "titulo": "Gestión de la Demanda",   "gridshape": "circular", "es_estructura": False},
    #         {"clave": "radar_sop",        "id_dim": "DIM-02",
    #          "titulo": "S&OP",                    "gridshape": "circular", "es_estructura": False},
    #         ...
    #     ],
    #     "layout_pptx": {
    #         "radar_principal"   : "radar_demanda",
    #         "radares_inferiores": ["radar_sop", ...],
    #     },
    # },
}


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 2 — BENCHMARKS POR SUBDIMENSIÓN
# Fuente: base histórica de diagnósticos Zonalogística.
# Estructura: "Nombre subdimensión": (puntaje_top, puntaje_media)
# ══════════════════════════════════════════════════════════════════════════════

BENCHMARKS = {
    # MOD-01 · Técnica y Operacional
    "Recibo"                    : (71.4, 66.3),
    "Almacenamiento"            : (72.0, 52.0),
    "Despacho"                  : (85.0, 60.7),
    "Control de inventarios"    : (86.7, 60.0),
    "Gestión de la planeación"  : (80.0, 60.0),
    # MOD-01 · Estructura
    "IT"                        : (80.0, 64.0),
    "Proceso"                   : (96.0, 60.0),
    "Estructura organizacional" : (100.0, 60.0),
    # MOD-01 · KPI
    "KPI Costo"                 : (90.0, 64.0),
    "KPI Servicio"              : (90.0, 54.3),
    "KPI Inventario"            : (90.0, 62.9),
    "KPI Seguridad"             : (90.0, 60.0),
    "Otros KPI"                 : (90.0, 40.0),
    # Cuando lleguen nuevos módulos se agregan aquí sus subdimensiones
}

# Orden de aparición de cada subdimensión en el radar general por módulo
ORDEN_RADAR_GENERAL = {
    "MOD-01": [
        "Recibo", "Almacenamiento", "Despacho", "Control de inventarios",
        "Gestión de la planeación",
        "IT", "Proceso", "Estructura organizacional",
        "KPI Costo", "KPI Servicio", "KPI Inventario", "KPI Seguridad", "Otros KPI",
    ],
}


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 3 — HELPERS INTERNOS
# ══════════════════════════════════════════════════════════════════════════════

def _codificar_logo(ruta: str | None) -> str | None:
    """Convierte un logo a base64 para incrustarlo en el HTML/PPTX."""
    if not ruta or not os.path.exists(ruta):
        if ruta:
            print(f"⚠️  Logo no encontrado en '{ruta}'. Se omite del encabezado.")
        return None
    with open(ruta, "rb") as f:
        datos = f.read()
    ext  = os.path.splitext(ruta)[1].lstrip(".").lower()
    mime = "png" if ext == "png" else ext
    return f"data:image/{mime};base64,{base64.b64encode(datos).decode('utf-8')}"


def _bench(subdimension: str) -> tuple[float, float]:
    """Retorna (puntaje_top, puntaje_media) para una subdimensión."""
    return BENCHMARKS.get(subdimension, (80.0, 60.0))


def _radar_desde_subdims(subdimensiones: dict) -> dict:
    """
    Construye el bloque de datos de un radar a partir de un dict
    {nombre_subdim: score_float}.
    """
    cats    = list(subdimensiones.keys())
    scores  = list(subdimensiones.values())
    tops    = [_bench(s)[0] for s in cats]
    medias  = [_bench(s)[1] for s in cats]
    return {
        "categorias"    : cats,
        "calificacion"  : [round(v, 1) for v in scores],
        "puntaje_top"   : tops,
        "puntaje_media" : medias,
    }


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 4 — CONSTRUCCIÓN DEL PAYLOAD
# ══════════════════════════════════════════════════════════════════════════════

def construir_payload(
    id_modulo     : str,
    empresa       : str,
    scores        : dict,
    ruta_logo     : str | None = None,
    ruta_logo_zona: str | None = None,
) -> dict:
    """
    Construye el diccionario DATOS_DIAGNOSTICO listo para inyectar en el HTML
    y generar el PPTX, a partir de los scores calculados por scoring.py.

    Parámetros
    ----------
    id_modulo      : str   → ej. "MOD-01"
    empresa        : str   → nombre de la empresa evaluada
    scores         : dict  → resultado de scoring.calcular_scores()
    ruta_logo      : str   → ruta al logo del cliente (opcional)
    ruta_logo_zona : str   → ruta al logo de Zonalogística (opcional)

    Retorna
    -------
    dict con la estructura completa que espera el dashboard_template.html
    """
    cfg        = CONFIG_MODULOS[id_modulo]
    dimensiones = scores["dimensiones"]

    # ── Gauges ────────────────────────────────────────────────────────────────
    gauges = [
        {
            "id"   : "gauge-general",
            "title": f"Madurez {cfg['centro_evaluado']}",
            "value": round(scores["score_general"], 1),
        }
    ]
    for rd in cfg["radares_dimensiones"]:
        id_dim = rd["id_dim"]
        if id_dim in dimensiones:
            gauges.append({
                "id"   : f"gauge-{rd['clave'].replace('radar_', '')}",
                "title": rd["titulo"],
                "value": round(dimensiones[id_dim]["score"], 1),
            })

    # ── Barras agrupadas por dimensión ────────────────────────────────────────
    nombres_dim, scores_dim, tops_dim, medias_dim = [], [], [], []
    for rd in cfg["radares_dimensiones"]:
        id_dim = rd["id_dim"]
        if id_dim not in dimensiones:
            continue
        data   = dimensiones[id_dim]
        subdims = data["subdimensiones"]
        nombres_dim.append(rd["titulo"])
        scores_dim.append(round(data["score"], 1))
        tops_dim.append(round(
            sum(_bench(s)[0] for s in subdims) / len(subdims), 1))
        medias_dim.append(round(
            sum(_bench(s)[1] for s in subdims) / len(subdims), 1))

    barras_dimensiones = {
        "categorias"    : nombres_dim,
        "calificacion"  : scores_dim,
        "puntaje_top"   : tops_dim,
        "puntaje_media" : medias_dim,
    }

    # ── Radar triangular (una punta por dimensión) ────────────────────────────
    radar_triangular = {
        "categorias"    : nombres_dim,
        "calificacion"  : scores_dim,
        "puntaje_top"   : tops_dim,
        "puntaje_media" : medias_dim,
    }

    # ── Radar general (todas las subdimensiones en orden definido) ────────────
    orden = ORDEN_RADAR_GENERAL.get(id_modulo, [])
    subdim_scores = {}
    for data in dimensiones.values():
        subdim_scores.update(data["subdimensiones"])

    cats_g, cal_g, top_g, med_g = [], [], [], []
    for subdim in orden:
        if subdim in subdim_scores:
            bench = _bench(subdim)
            cats_g.append(subdim)
            cal_g.append(round(subdim_scores[subdim], 1))
            top_g.append(bench[0])
            med_g.append(bench[1])

    radar_general = {
        "categorias"    : cats_g,
        "calificacion"  : cal_g,
        "puntaje_top"   : top_g,
        "puntaje_media" : med_g,
    }

    # ── Radares por dimensión (dinámicos) ─────────────────────────────────────
    # Lista estructurada — usada por generar_pptx() para iterar
    radares_dimensiones = []
    for rd in cfg["radares_dimensiones"]:
        id_dim = rd["id_dim"]
        if id_dim not in dimensiones:
            continue
        datos_radar = _radar_desde_subdims(dimensiones[id_dim]["subdimensiones"])
        radares_dimensiones.append({
            "clave"        : rd["clave"],
            "titulo"       : rd["titulo"],
            "gridshape"    : rd["gridshape"],
            "es_estructura": rd["es_estructura"],
            "datos"        : datos_radar,
        })

    # ── Ensamblaje del payload ────────────────────────────────────────────────
    payload = {
        "meta": {
            "empresa"         : empresa,
            "centro_evaluado" : cfg["centro_evaluado"],
            "logo_base64"     : _codificar_logo(ruta_logo),
            "logo_zona_base64": _codificar_logo(ruta_logo_zona),
        },
        "gauges"              : gauges,
        "barras_dimensiones"  : barras_dimensiones,
        "radar_general"       : radar_general,
        "radar_triangular"    : radar_triangular,
        # Lista dinámica — para generar_pptx() y módulos futuros
        "radares_dimensiones" : radares_dimensiones,
    }

    # ── Aliases de compatibilidad hacia atrás ─────────────────────────────────
    # El dashboard_template.html actual espera las claves originales por nombre.
    # Estos aliases permiten que el HTML siga funcionando sin ningún cambio.
    for rd in radares_dimensiones:
        payload[rd["clave"]] = rd["datos"]

    return payload


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 5 — VALIDACIÓN
# ══════════════════════════════════════════════════════════════════════════════

def validar_payload(payload: dict) -> bool:
    """
    Valida el payload antes de exportar.
    Detecta errores de longitud de series y valores fuera de rango.
    """
    errores = []

    # Gauges entre 0 y 100
    for g in payload.get("gauges", []):
        if not (0 <= g["value"] <= 100):
            errores.append(
                f"Gauge '{g['title']}': valor {g['value']} fuera de rango [0, 100]")

    # Barras: series de igual longitud que categorías
    bd    = payload.get("barras_dimensiones", {})
    n_cat = len(bd.get("categorias", []))
    for clave, vals in bd.items():
        if clave == "categorias":
            continue
        if len(vals) != n_cat:
            errores.append(
                f"barras_dimensiones['{clave}']: {len(vals)} valores, "
                f"se esperaban {n_cat}")

    # Radar general y triangular
    for nombre in ("radar_general", "radar_triangular"):
        r = payload.get(nombre, {})
        n = len(r.get("categorias", []))
        for clave, vals in r.items():
            if clave == "categorias":
                continue
            if len(vals) != n:
                errores.append(
                    f"{nombre}['{clave}']: {len(vals)} valores, "
                    f"se esperaban {n}")
        if "calificacion" not in r:
            errores.append(f"{nombre}: falta la serie 'calificacion'")

    # Radares por dimensión (lista dinámica)
    for rd in payload.get("radares_dimensiones", []):
        datos = rd.get("datos", {})
        n     = len(datos.get("categorias", []))
        for clave, vals in datos.items():
            if clave == "categorias":
                continue
            if len(vals) != n:
                errores.append(
                    f"radar '{rd['clave']}' → '{clave}': {len(vals)} valores, "
                    f"se esperaban {n}")
        if "calificacion" not in datos:
            errores.append(f"radar '{rd['clave']}': falta la serie 'calificacion'")

    if errores:
        print("⚠️  Inconsistencias en el payload:")
        for e in errores:
            print(f"   - {e}")
        return False

    print("✅ Validación correcta.")
    return True


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 6 — EXPORTACIÓN HTML
# ══════════════════════════════════════════════════════════════════════════════

def exportar_json(payload: dict, ruta_salida: str = "datos_diagnostico.json") -> None:
    """Exporta el payload como JSON (útil para depuración)."""
    with open(ruta_salida, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)
    print(f"✅ JSON generado: {ruta_salida}")


def generar_html(
    payload  : dict,
    template : str = "dashboard_template.html",
    salida   : str = "dashboard_output.html",
) -> None:
    """
    Inyecta DATOS_DIAGNOSTICO en el template HTML y guarda el archivo final.
    El bloque de datos se inserta justo después del cierre </style>,
    conservando el comportamiento original del template.
    """
    with open(template, "r", encoding="utf-8") as f:
        contenido = f.read()

    datos_js     = json.dumps(payload, ensure_ascii=False)
    bloque_datos = f"\n<script>\nconst DATOS_DIAGNOSTICO = {datos_js};\n</script>\n"
    contenido_final = contenido.replace("</style>", "</style>" + bloque_datos, 1)

    with open(salida, "w", encoding="utf-8") as f:
        f.write(contenido_final)
    print(f"✅ Dashboard HTML generado: {salida}")


def generar_html_bytes(
    payload : dict,
    template: str = "dashboard_template.html",
) -> bytes:
    """
    Igual que generar_html() pero retorna el HTML como bytes en memoria,
    sin escribir archivo en disco. Usado por Streamlit para el botón de descarga.
    """
    with open(template, "r", encoding="utf-8") as f:
        contenido = f.read()

    datos_js        = json.dumps(payload, ensure_ascii=False)
    bloque_datos    = f"\n<script>\nconst DATOS_DIAGNOSTICO = {datos_js};\n</script>\n"
    contenido_final = contenido.replace("</style>", "</style>" + bloque_datos, 1)
    return contenido_final.encode("utf-8")


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 7 — EXPORTACIÓN POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def _fig_to_png(fig, width: int, height: int):
    """Convierte una figura Plotly a bytes PNG en memoria."""
    import io
    buf = io.BytesIO()
    fig.write_image(buf, format="png", width=width, height=height, scale=2)
    buf.seek(0)
    return buf


def _construir_gauge(g: dict):
    """Reconstruye la figura Plotly de un gauge con la paleta corporativa."""
    import plotly.graph_objects as go
    FONT = "Poppins, Arial"
    fig  = go.Figure(go.Indicator(
        mode  = "gauge+number",
        value = g["value"],
        number= {"suffix": "%", "font": {"family": FONT, "size": 36, "color": "#003049"}},
        gauge = {
            "axis": {
                "range": [0, 100], "visible": True, "ticksuffix": "%",
                "tickfont": {"family": FONT, "size": 10, "color": "#9AA1AC"}, "dtick": 20,
            },
            "bar"    : {"color": "#003049", "thickness": 0.30},
            "bgcolor": "white", "borderwidth": 0,
            "steps"  : [
                {"range": [0,  64], "color": "#ff0303"},
                {"range": [64, 85], "color": "#ffcb03"},
                {"range": [85, 100],"color": "#9ACD00"},
            ],
        },
        title = {"text": g["title"], "font": {"family": FONT, "size": 13, "color": "#003049"}},
    ))
    fig.update_layout(
        margin        = {"t": 60, "b": 30, "l": 30, "r": 30},
        paper_bgcolor = "white",
        font          = {"family": FONT},
    )
    return fig


def _construir_barras(b: dict):
    """Reconstruye el gráfico de barras agrupadas con la paleta corporativa."""
    import plotly.graph_objects as go
    FONT = "Poppins, Arial"
    INK  = "#9AA1AC"
    fig  = go.Figure()
    for valores, color, nombre in [
        (b["puntaje_top"],  "#9ACD00", "Puntaje Top (cuartil superior)"),
        (b["puntaje_media"],"#0056A6", "Puntaje en la media"),
        (b["calificacion"], "#00A3E0", "Calificación"),
    ]:
        fig.add_trace(go.Bar(
            x=b["categorias"], y=valores, name=nombre, marker_color=color,
            text=[f"{v}%" for v in valores], textposition="outside",
            textfont={"family": FONT, "size": 11, "color": INK},
        ))
    fig.update_layout(
        barmode      = "group", bargap=0.22, bargroupgap=0.12,
        yaxis        = {"range": [0, 110], "ticksuffix": "%",
                        "tickfont": {"family": FONT, "size": 11, "color": INK},
                        "gridcolor": "#EEF0F3"},
        xaxis        = {"tickfont": {"family": FONT, "size": 11, "color": INK}},
        legend       = {"font": {"family": FONT, "size": 11, "color": INK},
                        "orientation": "h", "yanchor": "bottom", "y": -0.28,
                        "xanchor": "center", "x": 0.5},
        margin       = {"t": 50, "b": 80, "l": 50, "r": 20},
        paper_bgcolor= "white", plot_bgcolor="white",
        title        = {"text": "Calificación vs. Benchmark por Dimensión",
                        "font": {"family": FONT, "size": 14, "color": "#003049"}, "x": 0.5},
    )
    return fig


def _construir_radar(r: dict, titulo: str, gridshape: str = "circular",
                     es_estructura: bool = False):
    """
    Reconstruye cualquier radar Plotly con la paleta corporativa.
    Funciona para cualquier módulo: solo necesita el dict con
    categorias / calificacion / puntaje_top / puntaje_media.
    """
    import plotly.graph_objects as go
    FONT       = "Poppins, Arial"
    INK        = "#9AA1AC"
    C_TOP      = "#9ACD00"
    C_MED      = "#0056A6"
    C_CAL      = "#00A3E0"
    C_CAL_FILL = "rgba(0,163,224,0.10)"

    def cerrar(arr): return arr + [arr[0]]

    catC = cerrar(r["categorias"])
    trazos = [
        go.Scatterpolar(
            r=cerrar(r["puntaje_top"]),  theta=catC, fill="none",
            line={"color": C_TOP, "width": 3},
            marker={"size": 5 if not es_estructura else 6, "color": C_TOP},
            name="Puntaje Top",
        ),
        go.Scatterpolar(
            r=cerrar(r["puntaje_media"]), theta=catC, fill="none",
            line={"color": C_MED, "width": 3},
            marker={"size": 5 if not es_estructura else 6, "color": C_MED},
            name="Puntaje en la media",
        ),
        go.Scatterpolar(
            r=cerrar(r["calificacion"]),  theta=catC, fill="toself",
            fillcolor=C_CAL_FILL,
            line={"color": C_CAL, "width": 3},
            marker={"size": 6 if not es_estructura else 7, "color": C_CAL},
            name="Calificación",
        ),
    ]

    radial_cfg = {
        "visible"   : True, "range": [0, 100], "ticksuffix": "%",
        "tickfont"  : {"family": FONT, "size": 8 if es_estructura else 9, "color": INK},
    }
    if gridshape == "linear" or es_estructura:
        radial_cfg.update({"angle": 90, "tickangle": 0})

    fig = go.Figure(trazos)
    fig.update_layout(
        polar = {
            "gridshape"  : "linear" if es_estructura else gridshape,
            "radialaxis" : radial_cfg,
            "angularaxis": {"tickfont": {"family": FONT, "size": 10, "color": INK},
                            "linewidth": 0.5},
        },
        font      = {"family": FONT},
        showlegend= True,
        legend    = {"font": {"family": FONT, "size": 10, "color": INK},
                     "orientation": "h", "yanchor": "bottom", "y": -0.20,
                     "xanchor": "center", "x": 0.5},
        margin    = {"t": 50, "b": 80, "l": 60, "r": 60},
        paper_bgcolor= "white",
        title     = {"text": titulo,
                     "font": {"family": FONT, "size": 14, "color": "#003049"}, "x": 0.5},
    )
    return fig


def _slide_base(prs, titulo_seccion: str, ruta_logo: str | None,
                ruta_logo_zona: str | None):
    """Añade diapositiva con header navy, franja semáforo y logos."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    slide    = prs.slides.add_slide(prs.slide_layouts[6])
    W        = prs.slide_width
    H        = prs.slide_height
    HEADER_H = Inches(0.8)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF4, 0xF5, 0xF7)

    hdr = slide.shapes.add_shape(1, 0, 0, W, HEADER_H)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x49)
    hdr.line.fill.background()

    FRANJA_H = Emu(60000)
    for color_hex, x_frac in [("#ff0303", 0), ("#ffcb03", 1/3), ("#9ACD00", 2/3)]:
        r2 = int(color_hex[1:3], 16)
        g2 = int(color_hex[3:5], 16)
        b2 = int(color_hex[5:7], 16)
        fr = slide.shapes.add_shape(
            1, int(W * x_frac), HEADER_H - FRANJA_H, int(W / 3) + 2, FRANJA_H)
        fr.fill.solid()
        fr.fill.fore_color.rgb = RGBColor(r2, g2, b2)
        fr.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(9), Inches(0.55))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = titulo_seccion
    run.font.size  = Pt(18)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name  = "Calibri"

    logo_h = Inches(0.52)
    for ruta, lx in [(ruta_logo_zona, W - Inches(4.3)), (ruta_logo, W - Inches(2.1))]:
        if ruta and os.path.exists(ruta):
            slide.shapes.add_picture(ruta, lx, Inches(0.14), height=logo_h)

    return slide, W, H, HEADER_H


def _construir_presentacion_pptx(
    payload       : dict,
    id_modulo     : str,
    ruta_logo     : str | None = None,
    ruta_logo_zona: str | None = None,
    verbose       : bool = False,
):
    """
    Arma el objeto Presentation con las 3 diapositivas del diagnóstico.
    Diapositiva 3 se construye dinámicamente según el layout_pptx del
    módulo activo, sin claves hardcodeadas. Compartida por generar_pptx()
    (guarda a disco) y generar_pptx_bytes() (botón de descarga de Streamlit).

    Requiere: pip install python-pptx plotly kaleido
    Kaleido necesita Google Chrome: ejecuta `plotly_get_chrome` si no está instalado.
    """
    from pptx import Presentation
    from pptx.util import Inches

    if verbose:
        print("⏳ Generando PowerPoint...")

    cfg          = CONFIG_MODULOS[id_modulo]
    layout_pptx  = cfg["layout_pptx"]

    # Índice rápido: clave → bloque de datos del radar
    idx_radares  = {rd["clave"]: rd for rd in payload["radares_dimensiones"]}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    CONTENT_TOP = 0.95
    CONTENT_H   = 7.5 - CONTENT_TOP - 0.25

    # ── Diapositiva 1: Gauges ─────────────────────────────────────────────────
    slide1, W, H, _ = _slide_base(
        prs, "1 · Resumen ejecutivo de madurez", ruta_logo, ruta_logo_zona)
    gw  = Inches(3.1);  gh = Inches(2.6)
    gy  = Inches(CONTENT_TOP + (CONTENT_H - 2.6) / 2)
    gap = (W - gw * 4) / 5
    for i, g in enumerate(payload["gauges"]):
        buf = _fig_to_png(_construir_gauge(g), 600, 500)
        slide1.shapes.add_picture(buf, gap + i * (gw + gap), gy, width=gw, height=gh)
    if verbose:
        print("  ✓ Diapositiva 1 (Gauges)")

    # ── Diapositiva 2: Barras + Radar general ─────────────────────────────────
    slide2, W, H, _ = _slide_base(
        prs, "2 · Consolidado de indicadores", ruta_logo, ruta_logo_zona)
    half_w = (W - Inches(0.5)) / 2
    ch     = Inches(CONTENT_H - 0.1);  ct = Inches(CONTENT_TOP)
    slide2.shapes.add_picture(
        _fig_to_png(_construir_barras(payload["barras_dimensiones"]), 900, 600),
        Inches(0.15), ct, width=half_w, height=ch)
    slide2.shapes.add_picture(
        _fig_to_png(_construir_radar(
            payload["radar_general"], "Madurez — Vista General"), 700, 600),
        Inches(0.15) + half_w + Inches(0.2), ct, width=half_w, height=ch)
    if verbose:
        print("  ✓ Diapositiva 2 (Barras + Radar general)")

    # ── Diapositiva 3: Dimensiones (dinámico por módulo) ─────────────────────
    slide3, W, H, _ = _slide_base(
        prs, "3 · Dimensiones de la evaluación", ruta_logo, ruta_logo_zona)
    ct    = Inches(CONTENT_TOP);  ch = Inches(CONTENT_H)
    tri_w = Inches(4.8)

    # Radar triangular (izquierda)
    slide3.shapes.add_picture(
        _fig_to_png(_construir_radar(
            payload["radar_triangular"], "Madurez Logística", gridshape="linear"), 620, 700),
        Inches(0.1), ct, width=tri_w, height=ch)

    rx = Inches(0.1) + tri_w + Inches(0.15)
    rw = W - rx - Inches(0.1)

    # Radar principal (arriba a la derecha)
    clave_principal   = layout_pptx["radar_principal"]
    rd_principal      = idx_radares.get(clave_principal)
    if rd_principal:
        top_h = ch * 0.50
        slide3.shapes.add_picture(
            _fig_to_png(_construir_radar(
                rd_principal["datos"], rd_principal["titulo"],
                gridshape=rd_principal["gridshape"],
                es_estructura=rd_principal["es_estructura"]), 700, 420),
            rx, ct, width=rw, height=top_h)
    else:
        top_h = Inches(0)

    # Radares inferiores (abajo a la derecha — se reparten el espacio)
    claves_inf  = layout_pptx["radares_inferiores"]
    rds_inf     = [idx_radares[c] for c in claves_inf if c in idx_radares]
    bot_h       = ch - top_h - Inches(0.1)
    bot_y       = ct + top_h + Inches(0.1)
    n_inf       = len(rds_inf)
    if n_inf > 0:
        ancho_inf = rw / n_inf
        for j, rd_inf in enumerate(rds_inf):
            slide3.shapes.add_picture(
                _fig_to_png(_construir_radar(
                    rd_inf["datos"], rd_inf["titulo"],
                    gridshape=rd_inf["gridshape"],
                    es_estructura=rd_inf["es_estructura"]), 500, 380),
                rx + j * ancho_inf, bot_y,
                width=ancho_inf, height=bot_h)

    if verbose:
        print("  ✓ Diapositiva 3 (Dimensiones)")

    return prs


def generar_pptx(
    payload       : dict,
    id_modulo     : str,
    ruta_logo     : str | None = None,
    ruta_logo_zona: str | None = None,
    salida        : str = "presentacion_diagnostico.pptx",
) -> None:
    """Genera el PowerPoint con 3 diapositivas y lo guarda en disco."""
    prs = _construir_presentacion_pptx(
        payload, id_modulo, ruta_logo, ruta_logo_zona, verbose=True)
    prs.save(salida)
    print(f"✅ PowerPoint generado: {salida}")


def generar_pptx_bytes(
    payload       : dict,
    id_modulo     : str,
    ruta_logo     : str | None = None,
    ruta_logo_zona: str | None = None,
) -> bytes:
    """
    Igual que generar_pptx() pero retorna los bytes en memoria.
    Usado por Streamlit para el botón de descarga sin escribir disco.
    """
    import io
    prs = _construir_presentacion_pptx(
        payload, id_modulo, ruta_logo, ruta_logo_zona, verbose=False)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 8 — MATRIZ DE PRIORIZACIÓN
# Sin cambios respecto a v1.0.
# En el futuro esta función recibirá los datos calculados automáticamente
# desde las estrategias y sus atributos (impacto, urgencia, inversión).
# ══════════════════════════════════════════════════════════════════════════════

def _construir_figura_matriz(
    estrategias : list,
    impacto     : list,
    urgencia    : list,
    inversion   : list,
    descripciones: list = None,
):
    """
    Arma la figura Plotly de la matriz de priorización de estrategias. Los
    cuatro parámetros de lista deben tener la misma longitud. Si
    `descripciones` es None, se usa el label de cada estrategia como
    descripción en el tooltip. Compartida por generar_matriz() (guarda a
    disco) y generar_matriz_bytes() (botón de descarga de Streamlit).
    """
    import plotly.graph_objects as go
    from collections import defaultdict
    import math

    FONT       = "Poppins, Arial, sans-serif"
    COLOR_NAVY = "#003049"
    INK_SOFT   = "#9AA1AC"
    COLOR_SCALE = [[0.0, "#0056A6"], [0.5, "#00A3E0"], [1.0, "#9ACD00"]]

    if descripciones is None:
        descripciones = list(estrategias)
    desc_cortas = [
        (desc[:120] + "...") if len(desc) > 120 else desc
        for desc in descripciones
    ]

    # Jitter para separar puntos con coordenadas idénticas
    JITTER_RADIO = 0.45
    grupos = defaultdict(list)
    for i, (x, y) in enumerate(zip(urgencia, impacto)):
        grupos[(x, y)].append(i)

    urgencia_jit = list(urgencia)
    impacto_jit  = list(impacto)
    for (x, y), indices in grupos.items():
        n = len(indices)
        if n == 1:
            continue
        for k, idx in enumerate(indices):
            angulo = 2 * math.pi * k / n
            urgencia_jit[idx] = x + JITTER_RADIO * math.cos(angulo)
            impacto_jit[idx]  = y + JITTER_RADIO * math.sin(angulo)

    posiciones_texto = [
        "top center", "middle right", "bottom center", "middle left",
        "top right", "bottom right", "bottom left", "top left",
    ]
    text_positions_tmp = []
    for (x, y), indices in grupos.items():
        n = len(indices)
        for k, idx in enumerate(indices):
            pos = "top center" if n == 1 else posiciones_texto[k % len(posiciones_texto)]
            text_positions_tmp.append((idx, pos))
    text_positions = [pos for _, pos in sorted(text_positions_tmp, key=lambda t: t[0])]

    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=urgencia_jit, y=impacto_jit,
        mode="markers+text",
        text=estrategias, textposition=text_positions,
        textfont=dict(family=FONT, size=11, color=COLOR_NAVY),
        marker=dict(
            size=inversion, sizemode="area",
            sizeref=2.0 * max(inversion) / (60.0 ** 2), sizemin=1,
            color=inversion, colorscale=COLOR_SCALE,
            showscale=True, opacity=1,
            colorbar=dict(
                title=dict(text="Inversión", font=dict(family=FONT, size=12, color=COLOR_NAVY)),
                tickfont=dict(family=FONT, color=INK_SOFT),
            ),
            line=dict(width=0.5, color=COLOR_NAVY),
        ),
        customdata=list(zip(urgencia, impacto, desc_cortas)),
        hovertemplate=(
            "<b>Palanca %{text}</b><br>"
            "<br>"
            "%{customdata[2]}<br>"
            "<br>"
            "Urgencia: %{customdata[0]}<br>"
            "Impacto: %{customdata[1]}<br>"
            "Inversión: %{marker.size:.0f}<br>"
            "<extra></extra>"
        ),
    ))

    for x0, x1, y0, y1 in [(5.0, 5.0, 0, 12.0), (0, 12.0, 5.0, 5.0)]:
        fig.add_shape(type="line", x0=x0, x1=x1, y0=y0, y1=y1,
                      line=dict(color=INK_SOFT, width=1, dash="dash"))

    for x, y, texto in [
        (7.5, 9.5, "Estratégicas"),  (2.5, 9.5, "Oportunidades"),
        (2.5, 0.5, "Básicas"),       (7.5, 0.5, "Operativas"),
    ]:
        fig.add_annotation(x=x, y=y, text=texto, showarrow=False,
                           font=dict(family=FONT, size=16, color=INK_SOFT))

    fig.update_layout(
        title=dict(text="<b>Matriz de priorización de Palancas</b>",
                   font=dict(family=FONT, size=18, color=COLOR_NAVY),
                   x=0.5, xanchor="center"),
        xaxis=dict(title=dict(text="Urgencia", font=dict(family=FONT, size=13, color=INK_SOFT)),
                   range=[0, 12.0], tickfont=dict(family=FONT, color=INK_SOFT),
                   gridcolor="#EEF0F3", linecolor="#EEF0F3"),
        yaxis=dict(title=dict(text="Impacto", font=dict(family=FONT, size=13, color=INK_SOFT)),
                   range=[0, 12.0], tickfont=dict(family=FONT, color=INK_SOFT),
                   gridcolor="#EEF0F3", linecolor="#EEF0F3"),
        font=dict(family=FONT),
        height=550, width=700,
        plot_bgcolor="white", paper_bgcolor="white",
        hoverlabel=dict(
            bgcolor="#003049",
            bordercolor="#A8DC00",
            font=dict(
                family="Poppins, Arial",
                size=12,
                color="#ffffff"
            ),
            align="left",
            namelength=-1
        ),
    )

    return fig


def generar_matriz(
    estrategias : list,
    impacto     : list,
    urgencia    : list,
    inversion   : list,
    descripciones: list = None,
    salida      : str = "Matriz Priorización de Palancas.html",
) -> None:
    """
    Genera la matriz de priorización de estrategias como HTML interactivo
    y la guarda en disco.
    """
    fig = _construir_figura_matriz(estrategias, impacto, urgencia, inversion, descripciones)
    fig.write_html(salida)
    print(f"✅ Matriz de priorización generada: {salida}")


def generar_matriz_bytes(
    estrategias: list, impacto: list, urgencia: list, inversion: list,
    descripciones: list = None,
) -> bytes:
    """
    Igual que generar_matriz() pero retorna el HTML como bytes en memoria.
    Usado por Streamlit para el botón de descarga.
    """
    fig = _construir_figura_matriz(estrategias, impacto, urgencia, inversion, descripciones)
    return fig.to_html(full_html=True, include_plotlyjs="cdn").encode("utf-8")


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 8B — EXPORTACIÓN PDF DE ESTRATEGIAS
# Documento con la misma lista de estrategias recomendadas que se muestra en
# la pantalla de resultados (una tarjeta por estrategia, coloreada según su
# nivel de brecha).
# ══════════════════════════════════════════════════════════════════════════════

_BADGE_COLORES_PDF = {
    "crítica" : "#FF0303",
    "moderada": "#FFCB03",
    "leve"    : "#A8DC00",
}


def _construir_pdf_estrategias(
    estrategias   : list,
    empresa       : str,
    nombre_modulo : str,
    fecha         : str,
    ruta_logo_zona: str | None = None,
):
    """
    Arma el documento PDF (reportlab) con la lista de estrategias
    recomendadas. Compartida por generar_pdf_estrategias() (guarda a disco)
    y generar_pdf_estrategias_bytes() (botón de descarga de Streamlit).
    """
    import io
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, HRFlowable,
    )

    COLOR_NAVY = colors.HexColor("#003049")
    COLOR_INK  = colors.HexColor("#6B7280")

    styles = getSampleStyleSheet()
    estilo_titulo = ParagraphStyle(
        "ZL_Titulo", parent=styles["Title"], textColor=COLOR_NAVY, fontSize=18, spaceAfter=4)
    estilo_sub = ParagraphStyle(
        "ZL_Sub", parent=styles["Normal"], textColor=COLOR_INK, fontSize=10, spaceAfter=14)
    estilo_encabezado = ParagraphStyle(
        "ZL_Encabezado", parent=styles["Normal"], fontSize=9.5, spaceAfter=4)
    estilo_texto = ParagraphStyle(
        "ZL_Texto", parent=styles["Normal"], fontSize=10.5, leading=14,
        textColor=colors.HexColor("#1F2937"), spaceBefore=2, spaceAfter=4)
    estilo_meta = ParagraphStyle(
        "ZL_Meta", parent=styles["Normal"], fontSize=8.5, textColor=COLOR_INK)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter,
        topMargin=1.5 * cm, bottomMargin=1.5 * cm,
        leftMargin=1.8 * cm, rightMargin=1.8 * cm,
    )

    elementos = []

    if ruta_logo_zona and os.path.exists(ruta_logo_zona):
        elementos.append(Image(ruta_logo_zona, width=3.2 * cm, height=1.1 * cm, kind="proportional"))
        elementos.append(Spacer(1, 0.4 * cm))

    elementos.append(Paragraph("Estrategias recomendadas", estilo_titulo))
    elementos.append(Paragraph(f"{empresa} &middot; {nombre_modulo} &middot; {fecha}", estilo_sub))
    elementos.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#E5E7EB")))
    elementos.append(Spacer(1, 0.5 * cm))

    if not estrategias:
        elementos.append(Paragraph(
            "No se encontraron estrategias para las respuestas registradas.", estilo_texto))
    else:
        for est in estrategias:
            brecha    = est["nivel_brecha"]
            color_hex = _BADGE_COLORES_PDF.get(brecha.lower(), "#003049")

            encabezado = Paragraph(
                f'<b><font color="{color_hex}">{brecha}</font></b>'
                f'&nbsp;&nbsp;&nbsp;<font color="#6B7280">{est["subdimension"]}</font>',
                estilo_encabezado,
            )
            cuerpo = Paragraph(est["estrategia"], estilo_texto)
            meta = Paragraph(
                f"Impacto: {est['impacto']} &nbsp;&middot;&nbsp; Plazo: {est['plazo']}",
                estilo_meta,
            )

            tarjeta = Table([[[encabezado, cuerpo, meta]]], colWidths=[doc.width])
            tarjeta.setStyle(TableStyle([
                ("LINEBEFORE"  , (0, 0), (0, 0), 3, colors.HexColor(color_hex)),
                ("BACKGROUND"  , (0, 0), (0, 0), colors.white),
                ("LEFTPADDING" , (0, 0), (0, 0), 12),
                ("RIGHTPADDING", (0, 0), (0, 0), 10),
                ("TOPPADDING"  , (0, 0), (0, 0), 8),
                ("BOTTOMPADDING", (0, 0), (0, 0), 8),
            ]))
            elementos.append(tarjeta)
            elementos.append(Spacer(1, 0.3 * cm))

    doc.build(elementos)
    buf.seek(0)
    return buf


def generar_pdf_estrategias(
    estrategias   : list,
    empresa       : str,
    nombre_modulo : str,
    fecha         : str,
    ruta_logo_zona: str | None = None,
    salida        : str = "Estrategias.pdf",
) -> None:
    """Genera el PDF de estrategias recomendadas y lo guarda en disco."""
    buf = _construir_pdf_estrategias(estrategias, empresa, nombre_modulo, fecha, ruta_logo_zona)
    with open(salida, "wb") as f:
        f.write(buf.read())
    print(f"✅ PDF de estrategias generado: {salida}")


def generar_pdf_estrategias_bytes(
    estrategias   : list,
    empresa       : str,
    nombre_modulo : str,
    fecha         : str,
    ruta_logo_zona: str | None = None,
) -> bytes:
    """
    Igual que generar_pdf_estrategias() pero retorna los bytes en memoria.
    Usado por Streamlit para el botón de descarga sin escribir disco.
    """
    buf = _construir_pdf_estrategias(estrategias, empresa, nombre_modulo, fecha, ruta_logo_zona)
    return buf.read()


# ══════════════════════════════════════════════════════════════════════════════
# SECCIÓN 9 — EJECUCIÓN DIRECTA (modo script, sin Streamlit)
# Útil para pruebas o regenerar un diagnóstico desde consola.
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # Importar la capa de scoring para construir los datos desde respuestas reales
    import sys
    sys.path.insert(0, os.path.dirname(__file__))
    from utils.scoring import calcular_scores

    # ── Respuestas de ejemplo (equivalente a los datos de Haceb en v1.0) ──────
    respuestas_ejemplo = {
        "P-01": 33,  "P-02": 67,  "P-03": 33,  "P-04": 67,
        "P-05": 33,  "P-06": 0,   "P-07": 33,
        "P-08": 33,  "P-09": 33,
        "P-10": 33,  "P-11": 33,
        "P-12": 100, "P-13": 67,
        "P-14": 33,  "P-15": 67,
        "P-16": 67,  "P-17": 67,
        "P-18": 100, "P-19": 100,
        "P-20": 100, "P-21": 100,
    }

    scores  = calcular_scores("MOD-01", respuestas_ejemplo)
    payload = construir_payload(
        id_modulo      = "MOD-01",
        empresa        = "Haceb",
        scores         = scores,
        ruta_logo      = "haceb_logo.png",
        ruta_logo_zona = "zona_logo_final.png",
    )

    ok = validar_payload(payload)
    if ok:
        exportar_json(payload)
        generar_html(payload, template="assets/dashboard_template.html")
        generar_pptx(
            payload, id_modulo="MOD-01",
            ruta_logo="haceb_logo.png", ruta_logo_zona="zona_logo_final.png",
        )
        # Matriz: por ahora valores manuales (se automatizará con las estrategias)
        estrategias_labels = [str(i) for i in range(1, 16)]
        impacto_vals  = [10,10,10,10,10,10,7,7,10,10,10,10,8,7,10]
        urgencia_vals = [10,10,10,9,6,6,7,7,10,10,10,10,8,7,6]
        inversion_vals= [1,1,3,1,1,1,1,1,3,3,5,7,2,2,5]
        descripciones_ejemplo = [
            f"Descripción de ejemplo de la estrategia {i}"
            for i in estrategias_labels
        ]
        generar_matriz(estrategias_labels, impacto_vals, urgencia_vals,
                       inversion_vals, descripciones_ejemplo)

        print("\n✅ Todo listo:")
        print("   • dashboard_output.html                  → abre en el navegador")
        print("   • presentacion_diagnostico.pptx          → abre en PowerPoint")
        print("   • Matriz Priorización de Palancas.html   → abre en el navegador")
